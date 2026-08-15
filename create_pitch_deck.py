import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # CoalOmIT Exact Website Color Palette
    BG_DARK = RGBColor(6, 6, 16)          # #060610 (Deep Space Navy)
    CARD_BG = RGBColor(13, 13, 26)        # #0d0d1a (Dark Card Bg)
    CARD_BG_ALT = RGBColor(20, 20, 40)    # #141428 (Slightly lighter card)
    GREEN = RGBColor(0, 255, 136)         # #00ff88 (Neon Emerald Green)
    CYAN = RGBColor(0, 212, 255)          # #00d4ff (Electric Cyan)
    PURPLE = RGBColor(124, 58, 237)       # #7c3aed (Accent Purple)
    TEXT_WHITE = RGBColor(255, 255, 255)  # #ffffff (Primary White)
    TEXT_LIGHT = RGBColor(232, 232, 240)  # #e8e8f0 (Body Light Gray)
    MUTED = RGBColor(140, 145, 160)       # #8c91a0 (Muted Gray)
    BORDER_COLOR = RGBColor(30, 35, 55)   # Border subtle
    RED_HIGHLIGHT = RGBColor(255, 85, 85) # Vivid Warning Red

    def set_slide_background(slide):
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_DARK
        bg_shape.line.fill.background()
        return bg_shape

    def add_header_footer(slide, title_text, category_text="CoalOmIT — Green AI Platform"):
        set_slide_background(slide)

        # Header banner
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
        banner.fill.solid()
        banner.fill.fore_color.rgb = BG_DARK
        banner.line.color.rgb = GREEN
        banner.line.width = Pt(1.5)

        # Brand Badge "Cø"
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.25), Inches(0.65), Inches(0.65))
        badge.fill.solid()
        badge.fill.fore_color.rgb = GREEN
        badge.line.fill.background()
        tf_b = badge.text_frame
        tf_b.word_wrap = True
        p_b = tf_b.paragraphs[0]
        p_b.text = "Cø"
        p_b.font.bold = True
        p_b.font.size = Pt(22)
        p_b.font.color.rgb = BG_DARK
        p_b.alignment = PP_ALIGN.CENTER

        # Brand Name & Slide Title
        txBox = slide.shapes.add_textbox(Inches(1.3), Inches(0.15), Inches(11.5), Inches(0.85))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "CoalOmIT"
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = GREEN

        p2 = tf.add_paragraph()
        p2.text = title_text
        p2.font.bold = True
        p2.font.size = Pt(24)
        p2.font.color.rgb = TEXT_WHITE

        # Footer
        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4))
        footer.fill.solid()
        footer.fill.fore_color.rgb = BG_DARK
        footer.line.fill.background()
        tf_f = footer.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = category_text + "  |  Confidential Pitch Deck"
        p_f.font.size = Pt(10)
        p_f.font.color.rgb = MUTED
        p_f.alignment = PP_ALIGN.RIGHT

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)
    
    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(2.2))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = BG_DARK
    top_bar.line.color.rgb = GREEN
    top_bar.line.width = Pt(2)
    
    tx_logo = slide1.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.2))
    p_logo = tx_logo.text_frame.paragraphs[0]
    p_logo.text = "Cø CoalOmIT"
    p_logo.font.bold = True
    p_logo.font.size = Pt(48)
    p_logo.font.color.rgb = GREEN
    p_logo.alignment = PP_ALIGN.CENTER
    
    tx_title = slide1.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.333), Inches(2.8))
    tf_t = tx_title.text_frame
    p_t = tf_t.paragraphs[0]
    p_t.text = "Eliminate the Carbon Cost of Your AI Models"
    p_t.font.bold = True
    p_t.font.size = Pt(40)
    p_t.font.color.rgb = TEXT_WHITE
    p_t.alignment = PP_ALIGN.CENTER
    
    p_sub = tf_t.add_paragraph()
    p_sub.text = "CoalOmIT is an open-source toolkit that quantifies the energy and CO₂ trade-offs of AI model compression — in seconds, right inside your existing workflow."
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = TEXT_LIGHT
    p_sub.alignment = PP_ALIGN.CENTER
    
    bottom_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = BG_DARK
    bottom_bar.line.fill.background()
    tf_b1 = bottom_bar.text_frame
    p_b1 = tf_b1.paragraphs[0]
    p_b1.text = "Open Source · Apache 2.0 · Live on GitHub (AAYUSHI003/CoalOmit)"
    p_b1.font.size = Pt(11)
    p_b1.font.color.rgb = GREEN
    p_b1.alignment = PP_ALIGN.RIGHT

    # -------------------------------------------------------------
    # SLIDE 2: Team & Founders (Dark Theme 2-Founder Template)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_header_footer(slide2, "Founders & Leadership Team")
    
    positions = [Inches(1.5), Inches(7.0)]
    for i, x_pos in enumerate(positions):
        img_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.6), Inches(4.8), Inches(2.3))
        img_box.fill.solid()
        img_box.fill.fore_color.rgb = CARD_BG
        img_box.line.color.rgb = CYAN
        img_box.line.width = Pt(1)
        tf_img = img_box.text_frame
        p_i = tf_img.paragraphs[0]
        p_i.text = f"[ Insert Founder {i+1} Photo ]"
        p_i.font.size = Pt(14)
        p_i.font.color.rgb = MUTED
        p_i.alignment = PP_ALIGN.CENTER
        
        bio_box = slide2.shapes.add_textbox(x_pos, Inches(4.0), Inches(4.8), Inches(2.7))
        tf_bio = bio_box.text_frame
        tf_bio.word_wrap = True
        
        p_name = tf_bio.paragraphs[0]
        p_name.text = f"Founder {i+1} Name"
        p_name.font.bold = True
        p_name.font.size = Pt(18)
        p_name.font.color.rgb = TEXT_WHITE
        
        p_role = tf_bio.add_paragraph()
        p_role.text = "Co-Founder & Lead Engineer" if i == 0 else "Co-Founder & Systems Lead"
        p_role.font.bold = True
        p_role.font.size = Pt(13)
        p_role.font.color.rgb = GREEN
        
        p_desc = tf_bio.add_paragraph()
        p_desc.text = "• Undergraduate Student Engineer\n• Core Architecture & AI Tooling Lead" if i == 0 else "• Undergraduate Student Engineer\n• Systems Architecture & CI/CD Lead"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_LIGHT

    # -------------------------------------------------------------
    # SLIDE 3: Current Issues with AI Model Deployment (Dark Theme)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_header_footer(slide3, "Current Issues with AI Model Deployment")
    
    col1_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.5))
    col1_bg.fill.solid()
    col1_bg.fill.fore_color.rgb = CARD_BG
    col1_bg.line.color.rgb = BORDER_COLOR
    
    tf_c1 = col1_bg.text_frame
    tf_c1.word_wrap = True
    p1_h = tf_c1.paragraphs[0]
    p1_h.text = "ML Engineers & AI Teams"
    p1_h.font.bold = True
    p1_h.font.size = Pt(18)
    p1_h.font.color.rgb = CYAN
    p1_h.alignment = PP_ALIGN.CENTER
    
    bullets_ml = [
        "Completely blind to carbon delta when compressing models",
        "Optimize only for accuracy & latency, ignoring energy FLOPs",
        "No easy tool to measure per-inference kWh or regional CO₂",
        "Trial-and-error quantization (INT8 vs INT4) without trade-off data",
        "No automated PR checks for carbon bloat before merging code"
    ]
    for b in bullets_ml:
        p_b = tf_c1.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = TEXT_LIGHT
        
    col2_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.6), Inches(4.5))
    col2_bg.fill.solid()
    col2_bg.fill.fore_color.rgb = CARD_BG
    col2_bg.line.color.rgb = BORDER_COLOR
    
    tf_c2 = col2_bg.text_frame
    tf_c2.word_wrap = True
    p2_h = tf_c2.paragraphs[0]
    p2_h.text = "Enterprise & ESG Compliance"
    p2_h.font.bold = True
    p2_h.font.size = Pt(18)
    p2_h.font.color.rgb = CYAN
    p2_h.alignment = PP_ALIGN.CENTER
    
    bullets_esg = [
        "Silent accumulation of AI carbon bloat across cloud clusters",
        "Upcoming EU CSRD mandates corporate carbon reporting",
        "High electricity bills driven by unoptimized FP32 models",
        "Lack of bottom-up telemetry to prove green AI initiatives",
        "Delayed annual audit cycles instead of continuous CI/CD tracking"
    ]
    for b in bullets_esg:
        p_b = tf_c2.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = TEXT_LIGHT

    v_box = slide3.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.7))
    tf_v = v_box.text_frame
    p_v = tf_v.paragraphs[0]
    p_v.text = "ACCURACY & LATENCY ARE TRACKED DAILY — CARBON EMISSIONS REMAIN COMPLETELY BLIND"
    p_v.font.bold = True
    p_v.font.size = Pt(16)
    p_v.font.color.rgb = RED_HIGHLIGHT
    p_v.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 4: How It Works & Business Model (Dark Theme)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_header_footer(slide4, "How It Works & Business Model")
    
    tx_flow = slide4.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(6.5), Inches(4.5))
    tf_fl = tx_flow.text_frame
    tf_fl.word_wrap = True
    
    steps = [
        ("1. Input Model Script", "Developer runs `cac run model.py --methods int8,int4`"),
        ("2. FLOPs & Energy Calculation", "PyTorch native FlopCounterMode profiles hardware energy"),
        ("3. Regional Grid Mapping", "Energy mapped to monthly CO₂ across 35+ regional power grids"),
        ("4. Automated CI/CD Comment", "GitHub Action posts before/after carbon comparison on PRs")
    ]
    for step_title, step_desc in steps:
        p1 = tf_fl.add_paragraph() if tf_fl.paragraphs[0].text else tf_fl.paragraphs[0]
        p1.text = step_title
        p1.font.bold = True
        p1.font.size = Pt(14)
        p1.font.color.rgb = GREEN
        
        p2 = tf_fl.add_paragraph()
        p2.text = step_desc + "\n"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_LIGHT
        
    bm_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.5), Inches(4.7), Inches(4.4))
    bm_box.fill.solid()
    bm_box.fill.fore_color.rgb = CARD_BG
    bm_box.line.color.rgb = CYAN
    bm_box.line.width = Pt(1.5)
    
    tf_bm = bm_box.text_frame
    tf_bm.word_wrap = True
    p_bm_title = tf_bm.paragraphs[0]
    p_bm_title.text = "Business Model"
    p_bm_title.font.bold = True
    p_bm_title.font.size = Pt(22)
    p_bm_title.font.color.rgb = CYAN
    p_bm_title.alignment = PP_ALIGN.CENTER
    
    bm_points = [
        "Open-Source Core (Apache 2.0): Free developer CLI & PyTorch measurement engine.",
        "Enterprise SaaS Subscription: Hosted dashboard, historical carbon trends, and team analytics.",
        "Compliance Export API: EU CSRD & SEC ESG reporting exports billed per active inference node.",
        "Auto-Quantization Cloud: Automated Pareto compression engine for proprietary enterprise LLMs."
    ]
    for pt in bm_points:
        p_pt = tf_bm.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.size = Pt(11)
        p_pt.font.color.rgb = TEXT_LIGHT
        
    tag_box = slide4.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.7))
    tf_tag = tag_box.text_frame
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = "WE ARE BUILDING THE CARBON-AWARE COMPRESSION LAYER OVER EVERY ML PIPELINE"
    p_tag.font.bold = True
    p_tag.font.size = Pt(16)
    p_tag.font.color.rgb = RED_HIGHLIGHT
    p_tag.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 5: Core Architecture (Dark Theme)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_header_footer(slide5, "Our Technical Architecture & Capabilities")
    
    tx_arch = slide5.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(5.0))
    tf_a = tx_arch.text_frame
    tf_a.word_wrap = True
    
    arch_items = [
        "PyTorch 2.x Native FlopCounterMode: Accurately counts exact hardware FLOPs per inference without modifying model code.",
        "Multi-Method Quantization Engine: Supports FP32 Baseline, INT8 Dynamic, INT4 Weight-Only, AWQ, and GPTQ quantization backends.",
        "Regional Grid Carbon Intensity Database: Pre-computed intensity metrics for 35+ country grids (gCO₂/kWh) with live Electricity Maps API fallback.",
        "Automated CI/CD Bot (cac-action): GitHub Action container that posts detailed before/after carbon trade-off tables directly onto Pull Requests.",
        "Rich Developer CLI (cac-cli): Command-line tool producing color-coded terminal tables, Markdown reports, and JSON exports."
    ]
    for item in arch_items:
        title, desc = item.split(":", 1)
        p_t = tf_a.add_paragraph() if tf_a.paragraphs[0].text else tf_a.paragraphs[0]
        p_t.text = "• " + title + ":"
        p_t.font.bold = True
        p_t.font.size = Pt(15)
        p_t.font.color.rgb = GREEN
        
        p_d = tf_a.add_paragraph()
        p_d.text = "   " + desc + "\n"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_LIGHT

    # -------------------------------------------------------------
    # SLIDE 6: Unit Economics (Dark Theme Table)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_header_footer(slide6, "Unit Economics & Model Impact")
    
    rows, cols = 4, 4
    table_shape = slide6.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.5))
    table = table_shape.table
    
    headers = ["Compression Strategy", "Baseline (FP32)", "CoalOmIT (INT8 / INT4)", "Impact / Explanation"]
    col_widths = [Inches(3.2), Inches(2.2), Inches(2.6), Inches(3.733)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
        
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG_ALT
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.CENTER
        
    table_data = [
        ["Energy per 1k Inferences", "0.021 kWh", "0.009 kWh (-57%)", "57% energy reduction per batch via INT4 quantization"],
        ["Est. Monthly CO₂ (1M batch)", "128 kg CO₂", "55 kg CO₂ (-57%)", "Direct reduction of 73 kg CO₂ per model instance/mo"],
        ["Model Size & Latency", "440 MB / 42ms", "55 MB / 18ms (-57%)", "87% smaller footprint & 2.3x latency speedup"]
    ]
    for row_idx, row_content in enumerate(table_data, start=1):
        for col_idx, cell_value in enumerate(row_content):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if row_idx % 2 == 1 else BG_DARK
            cell.text = cell_value
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_WHITE if col_idx == 2 else TEXT_LIGHT
            if col_idx in [1, 2]:
                p.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 7: Differentiators (Dark Theme Boxes)
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_header_footer(slide7, "Our Key Differentiators")
    
    box_p = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(3.0), Inches(4.5))
    box_p.fill.solid()
    box_p.fill.fore_color.rgb = CARD_BG
    box_p.line.color.rgb = GREEN
    box_p.line.width = Pt(2)
    tf_p_title = box_p.text_frame
    p_pt = tf_p_title.paragraphs[0]
    p_pt.text = "\n\nPRODUCT"
    p_pt.font.bold = True
    p_pt.font.size = Pt(24)
    p_pt.font.color.rgb = GREEN
    p_pt.alignment = PP_ALIGN.CENTER
    
    tx_p_desc = slide7.shapes.add_textbox(Inches(4.1), Inches(1.6), Inches(8.4), Inches(2.2))
    tf_pd = tx_p_desc.text_frame
    tf_pd.word_wrap = True
    p_pd1 = tf_pd.paragraphs[0]
    p_pd1.text = "• PyTorch Native FLOP Profiling: No synthetic estimates; measures true hardware execution.\n• Regional Power Grid Database: Maps energy directly to local gCO₂/kWh in 35+ countries.\n• Zero-Code Developer Experience: Run single command `cac run model.py` on existing scripts."
    p_pd1.font.size = Pt(13)
    p_pd1.font.color.rgb = TEXT_LIGHT
    
    box_s = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.3), Inches(3.0), Inches(2.4))
    box_s.fill.solid()
    box_s.fill.fore_color.rgb = CARD_BG
    box_s.line.color.rgb = CYAN
    box_s.line.width = Pt(2)
    tf_s_title = box_s.text_frame
    p_st = tf_s_title.paragraphs[0]
    p_st.text = "\nSERVICE"
    p_st.font.bold = True
    p_st.font.size = Pt(24)
    p_st.font.color.rgb = CYAN
    p_st.alignment = PP_ALIGN.CENTER
    
    tx_s_desc = slide7.shapes.add_textbox(Inches(4.1), Inches(4.3), Inches(8.4), Inches(2.4))
    tf_sd = tx_s_desc.text_frame
    tf_sd.word_wrap = True
    p_sd1 = tf_sd.paragraphs[0]
    p_sd1.text = "• Automated CI/CD PR Bot: Posts trade-off comparison tables on every pull request.\n• EU CSRD Compliance Exports: Instant export of bottom-up audit logs for corporate ESG reports.\n• 100% Status & Pareto Visibility: Clear Pareto-optimal trade-offs across Accuracy, Latency, and CO₂."
    p_sd1.font.size = Pt(13)
    p_sd1.font.color.rgb = TEXT_LIGHT

    # -------------------------------------------------------------
    # SLIDE 8: Leveraging Technology (Dark Theme Table)
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_header_footer(slide8, "Leveraging Technology — Legacy vs CoalOmIT")
    
    r8, c8 = 6, 3
    t8_shape = slide8.shapes.add_table(r8, c8, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.0))
    t8 = t8_shape.table
    
    t8.columns[0].width = Inches(2.8)
    t8.columns[1].width = Inches(4.2)
    t8.columns[2].width = Inches(4.733)
    
    h8 = ["Capability", "Legacy / Manual Approach", "CoalOmIT Platform"]
    for j, h in enumerate(h8):
        cell = t8.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG_ALT
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.CENTER
        
    t8_data = [
        ["Carbon Measurement", "Static spreadsheet estimates & manual calculations", "Real-time PyTorch FLOPs & grid intensity lookup"],
        ["CI/CD Integration", "None — Post-deployment annual audits", "Automated GitHub Action PR comments (cac-action)"],
        ["Quantization Strategy", "Guesswork on INT8 vs FP16 accuracy drop", "Pareto trade-off matrix (Accuracy, Latency, CO₂)"],
        ["Enterprise ESG Exports", "Expensive annual sustainability consultants", "Instant CSRD & SEC compliant API data exports"],
        ["Execution Speed", "Weeks of manual auditing", "Instant CLI report in under 5 seconds"]
    ]
    for row_i, row_c in enumerate(t8_data, start=1):
        for col_j, val in enumerate(row_c):
            cell = t8.cell(row_i, col_j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if row_i % 2 == 1 else BG_DARK
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = GREEN if col_j == 2 else TEXT_LIGHT

    # -------------------------------------------------------------
    # SLIDE 9: Traction & Growth (Dark Theme)
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_slide_layout)
    add_header_footer(slide9, "Traction & Early Adoption")
    
    c9_1 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.8))
    c9_1.fill.solid()
    c9_1.fill.fore_color.rgb = CARD_BG
    c9_1.line.color.rgb = BORDER_COLOR
    tf_91 = c9_1.text_frame
    tf_91.word_wrap = True
    p91_h = tf_91.paragraphs[0]
    p91_h.text = "Developer & CLI Traction"
    p91_h.font.bold = True
    p91_h.font.size = Pt(18)
    p91_h.font.color.rgb = CYAN
    p91_h.alignment = PP_ALIGN.CENTER
    
    bullets_t1 = [
        "CLI Downloads growing ~25% month-on-month",
        "100+ GitHub Stars & growing open-source community",
        "Adopted by early AI startups & ML research labs",
        "Active PyTorch community contributors",
        "500+ Model compression reports generated"
    ]
    for b in bullets_t1:
        p_b = tf_91.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_LIGHT

    c9_2 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8))
    c9_2.fill.solid()
    c9_2.fill.fore_color.rgb = CARD_BG
    c9_2.line.color.rgb = BORDER_COLOR
    tf_92 = c9_2.text_frame
    tf_92.word_wrap = True
    p92_h = tf_92.paragraphs[0]
    p92_h.text = "Enterprise CI/CD Traction"
    p92_h.font.bold = True
    p92_h.font.size = Pt(18)
    p92_h.font.color.rgb = CYAN
    p92_h.alignment = PP_ALIGN.CENTER
    
    bullets_t2 = [
        "GitHub Action (cac-action) integrated in 20+ active PR pipelines",
        "Avg. 1,000+ monthly model inferences audited per enterprise node",
        "57% average carbon savings achieved across audited ML teams",
        "Strong interest from EU enterprise clients facing CSRD deadlines",
        "Partnering with green cloud infrastructure providers"
    ]
    for b in bullets_t2:
        p_b = tf_92.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_LIGHT

    # -------------------------------------------------------------
    # SLIDE 10: Usage Patterns & Impact Metrics (Dark Theme Grid)
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_slide_layout)
    add_header_footer(slide10, "Usage Patterns & Impact Metrics")
    
    grid = [
        (Inches(1.2), Inches(1.8), "57%", "Avg. CO₂ Reduction per Model", "Achieved via INT4 Weight-Only Quantization"),
        (Inches(7.0), Inches(1.8), "2.3x", "Inference Latency Speedup", "p50 latency improved from 42ms to 18ms"),
        (Inches(1.2), Inches(4.3), "87%", "Model Memory Savings", "Footprint reduced from 440MB to 55MB"),
        (Inches(7.0), Inches(4.3), "35+", "Regional Power Grids Mapped", "Live gCO₂/kWh grid intensity coverage")
    ]
    for x, y, num, label, subtext in grid:
        card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.1), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = GREEN
        card.line.width = Pt(1.5)
        
        tf_m = card.text_frame
        tf_m.word_wrap = True
        
        p_n = tf_m.paragraphs[0]
        p_n.text = num
        p_n.font.bold = True
        p_n.font.size = Pt(38)
        p_n.font.color.rgb = GREEN
        p_n.alignment = PP_ALIGN.CENTER
        
        p_l = tf_m.add_paragraph()
        p_l.text = label
        p_l.font.bold = True
        p_l.font.size = Pt(14)
        p_l.font.color.rgb = TEXT_WHITE
        p_l.alignment = PP_ALIGN.CENTER
        
        p_s = tf_m.add_paragraph()
        p_s.text = subtext
        p_s.font.size = Pt(11)
        p_s.font.color.rgb = MUTED
        p_s.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 11: Repeat Usage & Retention Cohorts (Dark Theme)
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_slide_layout)
    add_header_footer(slide11, "Repeat Usage & Pipeline Retention")
    
    r11, c11 = 5, 5
    t11_shape = slide11.shapes.add_table(r11, c11, Inches(0.8), Inches(1.6), Inches(11.733), Inches(2.5))
    t11 = t11_shape.table
    
    headers11 = ["Cohort Month", "Active Repos", "Month 1 Retention", "Month 2 Retention", "Month 3 Retention"]
    for j, h in enumerate(headers11):
        cell = t11.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG_ALT
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.CENTER
        
    data11 = [
        ["Q1 2026", "25 Repos", "100%", "84%", "76%"],
        ["Q2 2026", "65 Repos", "100%", "88%", "81%"],
        ["Q3 2026", "140 Repos", "100%", "91%", "—"],
        ["Q4 2026 (Est)", "300 Repos", "100%", "—", "—"]
    ]
    for r_i, row in enumerate(data11, start=1):
        for c_j, val in enumerate(row):
            cell = t11.cell(r_i, c_j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if r_i % 2 == 1 else BG_DARK
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_LIGHT
            p.alignment = PP_ALIGN.CENTER

    tx_ret = slide11.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.733), Inches(2.0))
    tf_r = tx_ret.text_frame
    
    p_r1 = tf_r.paragraphs[0]
    p_r1.text = "Average % of Repos retaining GitHub Action in 2nd month : 88%"
    p_r1.font.bold = True
    p_r1.font.size = Pt(18)
    p_r1.font.color.rgb = GREEN
    
    p_r2 = tf_r.add_paragraph()
    p_r2.text = "Average % of Repos retaining GitHub Action in 3rd month : 79%"
    p_r2.font.bold = True
    p_r2.font.size = Pt(18)
    p_r2.font.color.rgb = GREEN

    # -------------------------------------------------------------
    # SLIDE 12: Competitive Landscape (Dark Theme)
    # -------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_slide_layout)
    add_header_footer(slide12, "Competitive Landscape & Positioning")
    
    r12, c12 = 4, 3
    t12_shape = slide12.shapes.add_table(r12, c12, Inches(0.8), Inches(1.6), Inches(11.733), Inches(4.8))
    t12 = t12_shape.table
    
    t12.columns[0].width = Inches(3.0)
    t12.columns[1].width = Inches(4.2)
    t12.columns[2].width = Inches(4.533)
    
    h12 = ["Category", "Existing Competitors", "CoalOmIT Advantage"]
    for j, h in enumerate(h12):
        cell = t12.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG_ALT
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.CENTER
        
    data12 = [
        ["General Carbon Calculators\n(CodeCarbon, EcoLogITS)", "• Pure emissions tracking only\n• No model compression capabilities\n• Static top-level estimations", "• Combines carbon tracking + quantization\n• PyTorch native hardware FLOPs\n• Actionable CO₂ reduction engine"],
        ["ML Experiment Trackers\n(WandB, MLflow)", "• Focuses on accuracy & training loss\n• No built-in carbon grid intensity lookup\n• No automated PR carbon comments", "• Dedicated green AI metrics\n• Automated GitHub Action PR bot\n• Pareto trade-off matrix per model"],
        ["Enterprise ESG Platforms\n(Perspectium, Salesforce Net Zero)", "• Top-down annual estimates\n• No ML engineering visibility\n• High annual consulting costs", "• Bottom-up engineering telemetry\n• Instant API exports for EU CSRD\n• Developer-native workflow integration"]
    ]
    for r_i, row in enumerate(data12, start=1):
        for c_j, val in enumerate(row):
            cell = t12.cell(r_i, c_j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if r_i % 2 == 1 else BG_DARK
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = GREEN if col_j == 2 else TEXT_LIGHT

    # -------------------------------------------------------------
    # SLIDE 13: Roadmap (Dark Theme)
    # -------------------------------------------------------------
    slide13 = prs.slides.add_slide(blank_slide_layout)
    add_header_footer(slide13, "Product & Business Roadmap")
    
    r13, c13 = 5, 3
    t13_shape = slide13.shapes.add_table(r13, c13, Inches(0.8), Inches(1.6), Inches(11.733), Inches(4.8))
    t13 = t13_shape.table
    
    t13.columns[0].width = Inches(2.2)
    t13.columns[1].width = Inches(4.8)
    t13.columns[2].width = Inches(4.733)
    
    h13 = ["Timeline", "Product Milestones", "Business Milestones"]
    for j, h in enumerate(h13):
        cell = t13.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG_ALT
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.CENTER
        
    data13 = [
        ["Q1 2026", "• Launch `cac-core` PyTorch FLOP counter\n• Map 35+ country power grid intensity datasets\n• Build `cac-cli` rich terminal table output", "• Launch Apache 2.0 Open-Source repository\n• Build developer community & initial contributors"],
        ["Q2 2026", "• Release `cac-action` GitHub Action PR bot\n• Support INT8 Dynamic & INT4 Weight-Only\n• Add Markdown & JSON report export formats", "• Onboard top 20 open-source ML repositories\n• Initiate Enterprise Beta partner program"],
        ["Q3 2026", "• Launch Cloud Hosted API for telemetry\n• Add live Electricity Maps API fallback\n• Develop EU CSRD compliance export module", "• Launch Enterprise Paid Tier ($/node/mo)\n• Establish sales pipeline in EU & North America"],
        ["Q4 2026", "• Auto-Quantization Engine for LLMs\n• Multi-cloud cluster energy monitoring (AWS/GCP)\n• Real-time cloud cost vs CO₂ optimization", "• Scale to 100+ paying enterprise ML teams\n• Expand team in AI engineering & ESG sales"]
    ]
    for r_i, row in enumerate(data13, start=1):
        for c_j, val in enumerate(row):
            cell = t13.cell(r_i, c_j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if r_i % 2 == 1 else BG_DARK
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_LIGHT

    output_path = "c:\\Users\\hp\\CoalOmIT\\CoalOmIT_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully at: {output_path}")

if __name__ == "__main__":
    create_deck()
